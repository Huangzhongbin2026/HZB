from __future__ import annotations

from io import BytesIO
from pathlib import Path

from docx import Document
from fastapi import HTTPException
from pypdf import PdfReader
from pptx import Presentation


class DocumentIngestService:
    supported_suffixes = {".txt", ".md", ".docx", ".pptx", ".pdf"}

    def extract_text(self, filename: str, content: bytes) -> tuple[str, str]:
        suffix = Path(filename).suffix.lower()
        if suffix not in self.supported_suffixes:
            raise HTTPException(status_code=400, detail="unsupported document type")

        if suffix in {".txt", ".md"}:
            text = content.decode("utf-8", errors="ignore")
        elif suffix == ".docx":
            document = Document(BytesIO(content))
            text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
        elif suffix == ".pptx":
            presentation = Presentation(BytesIO(content))
            chunks: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and str(shape.text).strip():
                        chunks.append(str(shape.text).strip())
            text = "\n".join(chunks)
        else:
            reader = PdfReader(BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)

        normalized = "\n".join(line.rstrip() for line in text.splitlines())
        normalized = "\n".join(line for line in normalized.splitlines() if line.strip())
        title = Path(filename).stem or "uploaded-document"
        return title, normalized.strip()


document_ingest_service = DocumentIngestService()